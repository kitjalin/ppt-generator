
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import os
from datetime import datetime

# 初始化
st.title("安衛環稽核異常處理簡報自動產生器")

# 下拉選單選項
companies = ["A公司", "B公司", "C公司"]
departments = ["總務部", "品保部", "生產部"]
handlers = ["王小明", "陳美麗", "張三"]

company = st.selectbox("選擇公司", companies)
department = st.selectbox("選擇部門", departments)
handler = st.selectbox("處理人員", handlers)
date = st.date_input("異常日期", datetime.today())
description = st.text_area("異常描述")
photos = st.file_uploader("上傳照片 (可多張)", type=["jpg", "png"], accept_multiple_files=True)

if st.button("產生簡報"):
    ppt = Presentation("template.pptx")
    slide_layout = ppt.slide_layouts[1]

    slide = ppt.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = f"{company} - {department}"
    content.text = f"""
異常日期: {date.strftime('%Y-%m-%d')}
處理人員: {handler}
異常描述: {description}
"""

    for photo in photos:
        img_path = f"temp_{photo.name}"
        with open(img_path, "wb") as f:
            f.write(photo.read())
        ppt.slides.add_slide(ppt.slide_layouts[6]).shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(5))
        os.remove(img_path)

    out_pptx = "異常報告簡報.pptx"
    ppt.save(out_pptx)

    with open(out_pptx, "rb") as f:
        st.download_button("下載簡報", f, file_name=out_pptx)
